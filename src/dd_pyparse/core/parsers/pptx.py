from io import BytesIO
from pathlib import Path

import pptx
from loguru import logger
from PIL import UnidentifiedImageError
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_MEDIA_TYPE

from dd_pyparse.core.parsers.base import FileParser, get_file_meta
from dd_pyparse.core.parsers.image import ImageParser
from dd_pyparse.core.utils.general import safe_open
from dd_pyparse.core.utils.msoffice import convert_table_to_text, parse_meta
from dd_pyparse.schemas.data.image import Image


class PptxParser(FileParser):
    @staticmethod
    def order_shapes(shapes):
        """Order shapes by their position"""
        return sorted(shapes, key=lambda x: (x.top or 0, x.left or 0))

    @staticmethod
    def parse_children(
        presentation: pptx.presentation.Presentation,
        extract_children: bool = False,
        out_dir: Path = None,
        delimiter: str = " ",
        max_pages: int = None,
        notes_delimiter: str = "\nNotes: ",
        page_delimiter: str = "",
    ):
        """Parse children from presentation"""

        children = []
        out = {"num_pages": len(presentation.slides)}
        text = ""
        for i, slide in enumerate(presentation.slides):
            if max_pages is not None and i >= max_pages:
                break

            for shape in PptxParser.order_shapes(slide.shapes):
                if shape.has_text_frame:
                    new_text = PptxParser.parse_text(shape, delimiter=delimiter)
                    if new_text:
                        text = f"{text}{page_delimiter}{new_text}"
                elif shape.has_table:
                    table_text = PptxParser.parse_table(shape)
                    if table_text:
                        text = f"{text}{page_delimiter}{table_text}"
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    child = PptxParser.parse_image(shape, extract_children=extract_children, out_dir=out_dir)
                    if child:
                        children.append(child)
                elif shape.media_type == PP_MEDIA_TYPE.AUDIO:
                    logger.error("Found a audio file, but we don't support parsing it yet")
                    continue
                elif shape.media_type == PP_MEDIA_TYPE.MOVIE:
                    # save the movie file
                    logger.error("Found a video file, but we don't support parsing it yet")
                    continue
                else:
                    continue

            notes = PptxParser.parse_notes(slide)
            if notes:
                text = f"{text}{notes_delimiter}{notes}{page_delimiter}"

        if text:
            out["text"] = {"source": text.strip()}

        if children:
            out["children"] = children

        return out

    @staticmethod
    def parse_image(shape, extract_children: bool = False, min_resolution: int = -1, out_dir: Path = None) -> Image:
        """Get image from shape"""
        file_name = shape.image.filename
        file_ext = Path(file_name).suffix
        if file_ext.lower() == ".emf":
            return None

        child = get_file_meta(shape.image.blob)
        child["file_name"] = file_name
        child["file_extension"] = file_ext
        
        try:
            child |= ImageParser.parse(file=shape.image.blob)
        except UnidentifiedImageError:
            logger.warning(f"Failed to parse image {file_name}")

        if extract_children and all(dim > min_resolution for dim in [child["height"], child["width"]]):
            out_name = f"{child['hash']['md5']}{child.get('file_extension', '.bin')}"
            out_path = out_dir / out_name
            with safe_open(out_path, "wb") as fb:
                fb.write(shape.image.blob)
            child["absolute_path"] = out_path.absolute()
            logger.info(f"Extracted image {out_path}")

        return Image(**child)

    @staticmethod
    def parse_notes(slide):
        """Parse notes from shape"""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            return slide.notes_slide.notes_text_frame.text

    @staticmethod
    def parse_table(shape, fmt: str = "plain"):
        """Parse a table as text"""
        return convert_table_to_text(shape.table, fmt=fmt)

    @staticmethod
    def parse_text(shape, delimiter: str = " "):
        return delimiter.join([para.text for para in shape.text_frame.paragraphs])

    @staticmethod
    def parse(
        file: Path | bytes,
        extract_children: bool = False,
        out_dir: Path = None,
        delimiter: str = " ",
        max_pages: int = None,
        notes_delimiter: str = "\nNotes: ",
        page_delimiter: str = "",
        **kwargs,
    ) -> dict:
        """Parse a pptx file"""
        out = {}
        if isinstance(file, bytes):
            file = BytesIO(file)
        elif isinstance(file, (Path, str)):
            file = Path(file)
            out["absolute_path"] = file.absolute()
        else:
            raise TypeError(f"Cannot parse {type(file)}")

        presentation = pptx.Presentation(file)
        return {
            **parse_meta(presentation),
            **PptxParser.parse_children(
                presentation=presentation,
                extract_children=extract_children,
                out_dir=out_dir,
                delimiter=delimiter,
                max_pages=max_pages,
                notes_delimiter=notes_delimiter,
                page_delimiter=page_delimiter,
            ),
        }
